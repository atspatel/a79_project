import io
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


from utils.pexels_utils import get_image_stream_from_url


def apply_theme(ppt, theme_config):
    """Applies the theme configuration to all slides in the PowerPoint presentation."""

    for slide in ppt.slides:
        # Apply background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(
            *theme_config["colors"]["background_color"]
        )

        # Apply title font, font size, and color
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text_frame.paragraphs[0].font.name = theme_config["fonts"][
                "title_font"
            ]
            title_shape.text_frame.paragraphs[0].font.size = Pt(
                theme_config["font_sizes"]["title_size"]
            )
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                *theme_config["colors"]["title_color"]
            )

        # Apply content font, font size, and color to all content placeholders
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Skip title shape, it's already handled
                if shape == slide.shapes.title:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = theme_config["fonts"]["content_font"]
                    paragraph.font.size = Pt(theme_config["font_sizes"]["content_size"])
                    paragraph.font.color.rgb = RGBColor(
                        *theme_config["colors"]["content_color"]
                    )


def create_slide(ppt, slide_data):
    slide_layout = ppt.slide_layouts[slide_data["layout_id"]]
    slide = ppt.slides.add_slide(slide_layout)

    for placeholder in slide.placeholders:
        placeholder_name = placeholder.name
        matched_content = next(
            (x for x in slide_data["content"] if x["name"] == placeholder_name), None
        )
        if value := matched_content["value"]:
            if "picture" in placeholder_name.lower() and isinstance(value, str):
                try:
                    if image_url := matched_content["image_url"]:
                        placeholder.insert_picture(get_image_stream_from_url(image_url))
                except Exception as e:
                    print(f"Error inserting picture: {e}")

            elif isinstance(value, list):
                for line in value:
                    p = placeholder.text_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)  # Adjust font size if needed
            else:
                placeholder.text = value
    return slide


def generate_pptx_strem(slides, theme):
    ppt = Presentation()
    for slide_data in slides:
        create_slide(ppt, slide_data)

    if theme:
        apply_theme(ppt, theme)

    pptx_stream = io.BytesIO()
    ppt.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream


def save_ppt(ppt, output_file="presentation.pptx"):
    ppt.save(output_file)
    print("PowerPoint presentation saved as presentation.pptx")
